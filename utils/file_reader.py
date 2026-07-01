"""
📄 قارئ الملفات الخارق - v2.0
بيقرأ أي نوع ملف ويرجّع الصفحات منفصلة (للتقسيم)
"""
import os
import sys
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import PyPDF2
import pytesseract
from PIL import Image
from config import Config


class FileReader:
    """
    قارئ ملفات ذكي - بيقرأ أي نوع ملف ويستخرج النص منه

    الأنواع المدعومة:
    - PDF (نص + صور)
    - صور (PNG, JPG, BMP, GIF, WebP)
    - Word (DOCX)
    - PowerPoint (PPTX)
    - Excel (XLSX, CSV)
    - نصوص (TXT)
    """

    def __init__(self):
        """تهيئة قارئ الملفات"""
        # إعداد Tesseract لـ Windows
        if os.name == 'nt':
            tesseract_paths = [
                r'C:\Program Files\Tesseract-OCR\tesseract.exe',
                r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
                r'C:\Users\CS\AppData\Local\Tesseract-OCR\tesseract.exe',
            ]
            for path in tesseract_paths:
                if os.path.exists(path):
                    pytesseract.pytesseract.tesseract_cmd = path
                    break

    def read(self, uploaded_file):
        """
        القراءة الرئيسية - بتحدد نوع الملف وتقرأه

        Returns:
            dict: {
                'success': True/False,
                'text': النص المستخرج (كل الصفحات),
                'pages': list من الصفحات منفصلة,
                'total_pages': عدد الصفحات,
                'file_info': معلومات الملف,
                'error': رسالة الخطأ
            }
        """
        try:
            file_name = uploaded_file.name
            file_ext = file_name.split('.')[-1].lower()
            file_size = uploaded_file.size

            file_info = {
                'name': file_name,
                'extension': file_ext,
                'size': file_size,
                'size_readable': self._format_size(file_size)
            }

            # التحقق من الحجم
            if file_size > Config.MAX_FILE_SIZE:
                return self._error_response(
                    file_info,
                    f"❌ الملف كبير أوي! الحد الأقصى {self._format_size(Config.MAX_FILE_SIZE)}"
                )

            # التحقق من النوع
            if file_ext not in Config.ALL_EXTENSIONS:
                return self._error_response(
                    file_info,
                    f"❌ نوع الملف '.{file_ext}' مش مدعوم"
                )

            uploaded_file.seek(0)

            # القراءة حسب النوع
            readers = {
                'pdf': self._read_pdf,
                'txt': self._read_txt,
                'docx': self._read_docx,
                'doc': self._read_docx,
                'pptx': self._read_pptx,
                'ppt': self._read_pptx,
                'xlsx': self._read_excel,
                'xls': self._read_excel,
                'csv': self._read_csv,
                'png': self._read_image,
                'jpg': self._read_image,
                'jpeg': self._read_image,
                'bmp': self._read_image,
                'gif': self._read_image,
                'webp': self._read_image,
            }

            reader = readers.get(file_ext)
            if not reader:
                return self._error_response(
                    file_info,
                    f"❌ مفيش قارئ لنوع '.{file_ext}'"
                )

            # القراءة - بترجع pages (list)
            result = reader(uploaded_file)

            # لو الـ reader رجّع dict (PDF/PPTX اللي بيرجعوا صفحات)
            if isinstance(result, dict):
                pages = result.get('pages', [])
                text = result.get('text', '')
            else:
                # الباقي (txt, docx, image...) بيرجّع نص واحد
                pages = [result] if result and result.strip() else []
                text = result

            if not pages or not text or not text.strip():
                return self._error_response(
                    file_info,
                    "⚠️ مقدرتش أستخرج نص من الملف ده. ممكن يكون فاضي أو محمي."
                )

            total_pages = len(pages)
            file_info['total_pages'] = total_pages

            return {
                'success': True,
                'text': text.strip(),
                'pages': pages,
                'total_pages': total_pages,
                'file_info': file_info,
                'word_count': len(text.split()),
                'char_count': len(text),
                'error': None
            }

        except Exception as e:
            return self._error_response(
                {},
                f"❌ خطأ غير متوقع: {str(e)}"
            )

    def _error_response(self, file_info, error_msg):
        """رد موحّد للأخطاء"""
        return {
            'success': False,
            'text': '',
            'pages': [],
            'total_pages': 0,
            'file_info': file_info,
            'error': error_msg
        }

    # ============================================
    # 📕 قراءة PDF (مع كل الصفحات منفصلة)
    # ============================================

    def _read_pdf(self, file):
        """قراءة ملف PDF - يرجّع كل صفحة منفصلة"""
        try:
            reader = PyPDF2.PdfReader(file)
            total_pages = len(reader.pages)
            pages = []
            full_text = ""

            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    formatted_page = f"\n{'='*50}\n📄 صفحة {i+1} من {total_pages}\n{'='*50}\n{page_text}\n"
                    pages.append(page_text.strip())
                    full_text += formatted_page
                else:
                    # حتى لو الصفحة فاضية، نضيفها عشان الترقيم
                    pages.append("")

            # لو مفيش نص (PDF صور)
            if not full_text.strip():
                ocr_result = self._read_pdf_with_ocr(file)
                if isinstance(ocr_result, dict):
                    return ocr_result
                else:
                    # OCR رجع نص واحد، نقسمه
                    return {
                        'text': ocr_result,
                        'pages': [ocr_result] if ocr_result else []
                    }

            return {
                'text': full_text,
                'pages': pages
            }

        except Exception as e:
            return {
                'text': f"خطأ في قراءة PDF: {str(e)}",
                'pages': []
            }

    def _read_pdf_with_ocr(self, file):
        """قراءة PDF صور باستخدام OCR - مع الصفحات منفصلة"""
        try:
            import fitz  # PyMuPDF
            file.seek(0)
            doc = fitz.open(stream=file.read(), filetype="pdf")
            pages = []
            full_text = ""

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=200)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                try:
                    page_text = pytesseract.image_to_string(img, lang='ara+eng')
                except Exception:
                    page_text = pytesseract.image_to_string(img, lang='eng')

                page_text = page_text.strip() if page_text else ""
                pages.append(page_text)

                if page_text:
                    full_text += f"\n📄 صفحة {page_num + 1}:\n{page_text}\n"

            doc.close()
            return {
                'text': full_text,
                'pages': pages
            }

        except ImportError:
            return {
                'text': "⚠️ الملف ده PDF صور. ثبّت PyMuPDF: pip install PyMuPDF",
                'pages': []
            }
        except Exception as e:
            return {
                'text': f"خطأ في OCR: {str(e)}",
                'pages': []
            }

    # ============================================
    # 📝 قراءة نص عادي
    # ============================================

    def _read_txt(self, file):
        """قراءة ملف نصي"""
        encodings = ['utf-8', 'windows-1256', 'iso-8859-6', 'cp1252', 'latin-1']

        for encoding in encodings:
            try:
                file.seek(0)
                return file.read().decode(encoding)
            except (UnicodeDecodeError, AttributeError):
                continue

        file.seek(0)
        return file.read().decode('utf-8', errors='ignore')

    # ============================================
    # 📘 قراءة Word
    # ============================================

    def _read_docx(self, file):
        """قراءة ملف Word"""
        try:
            from docx import Document
            doc = Document(file)
            text = ""

            for para in doc.paragraphs:
                if para.text.strip():
                    if para.style.name.startswith('Heading'):
                        level = para.style.name[-1] if para.style.name[-1].isdigit() else '1'
                        text += f"\n{'#' * int(level)} {para.text}\n"
                    else:
                        text += para.text + "\n"

            for table in doc.tables:
                text += "\n📊 جدول:\n"
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    text += " | ".join(cells) + "\n"
                text += "\n"

            return text

        except Exception as e:
            return f"خطأ في قراءة Word: {str(e)}"

    # ============================================
    # 📊 قراءة PowerPoint (كل شريحة = صفحة)
    # ============================================

    def _read_pptx(self, file):
        """قراءة ملف PowerPoint - كل شريحة كصفحة منفصلة"""
        try:
            from pptx import Presentation
            prs = Presentation(file)
            pages = []
            full_text = ""

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = ""

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_text += paragraph.text + "\n"

                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_text += " | ".join(cells) + "\n"

                pages.append(slide_text.strip())

                full_text += f"\n{'='*40}\n🎞️ شريحة {slide_num}\n{'='*40}\n{slide_text}\n"

            return {
                'text': full_text,
                'pages': pages
            }

        except Exception as e:
            return {
                'text': f"خطأ في قراءة PowerPoint: {str(e)}",
                'pages': []
            }

    # ============================================
    # 📈 قراءة Excel
    # ============================================

    def _read_excel(self, file):
        """قراءة ملف Excel"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file, data_only=True)
            text = ""

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text += f"\n📊 شيت: {sheet_name}\n"
                text += "─" * 40 + "\n"

                for row in ws.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(cells):
                        text += " | ".join(cells) + "\n"

                text += "\n"

            return text

        except Exception as e:
            return f"خطأ في قراءة Excel: {str(e)}"

    # ============================================
    # 📄 قراءة CSV
    # ============================================

    def _read_csv(self, file):
        """قراءة ملف CSV"""
        try:
            import csv
            import io

            content = file.read().decode('utf-8', errors='ignore')
            reader = csv.reader(io.StringIO(content))

            text = "📊 جدول بيانات:\n"
            text += "─" * 40 + "\n"

            for row in reader:
                text += " | ".join(row) + "\n"

            return text

        except Exception as e:
            return f"خطأ في قراءة CSV: {str(e)}"

    # ============================================
    # 🖼️ قراءة الصور (OCR)
    # ============================================

    def _read_image(self, file):
        """قراءة النص من صورة"""
        try:
            image = Image.open(file)

            if image.mode != 'RGB':
                image = image.convert('RGB')

            width, height = image.size
            if width < 1000:
                ratio = 1000 / width
                image = image.resize((int(width * ratio), int(height * ratio)))

            text = ""
            try:
                text = pytesseract.image_to_string(image, lang='ara+eng')
            except Exception:
                try:
                    text = pytesseract.image_to_string(image, lang='eng')
                except Exception:
                    text = pytesseract.image_to_string(image)

            return text if text.strip() else "⚠️ مقدرتش أقرأ نص من الصورة. تأكد إن الصورة واضحة."

        except Exception as e:
            return f"خطأ في قراءة الصورة: {str(e)}"

    # ============================================
    # 🆕 دوال جديدة للتقسيم
    # ============================================

    @staticmethod
    def get_pages_range(pages, start_page, end_page):
        """
        الحصول على نص من نطاق صفحات معين

        Args:
            pages: list من الصفحات
            start_page: رقم البداية (يبدأ من 1)
            end_page: رقم النهاية (شامل)

        Returns:
            str: النص المجمّع من الصفحات
        """
        if not pages:
            return ""

        # تحويل من 1-based لـ 0-based
        start_idx = max(0, start_page - 1)
        end_idx = min(len(pages), end_page)

        selected_pages = pages[start_idx:end_idx]

        # دمج الصفحات
        combined_text = ""
        for i, page_text in enumerate(selected_pages, start=start_page):
            if page_text and page_text.strip():
                combined_text += f"\n{'='*50}\n📄 صفحة {i}\n{'='*50}\n{page_text}\n"

        return combined_text.strip()

    @staticmethod
    def suggest_ranges(total_pages, chunk_size=15):
        """
        اقتراح نطاقات للتقسيم

        Args:
            total_pages: عدد الصفحات الكلي
            chunk_size: حجم كل جزء (افتراضي 15)

        Returns:
            list: [(start, end), (start, end), ...]
        """
        if total_pages <= chunk_size:
            return [(1, total_pages)]

        ranges = []
        for start in range(1, total_pages + 1, chunk_size):
            end = min(start + chunk_size - 1, total_pages)
            ranges.append((start, end))

        return ranges

    # ============================================
    # 🔧 أدوات مساعدة
    # ============================================

    @staticmethod
    def _format_size(size_bytes):
        """تحويل حجم الملف لصيغة مقروءة"""
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.1f} KB"
        else:
            return f"{size_bytes / (1024 * 1024):.1f} MB"

    @staticmethod
    def get_file_icon(extension):
        """إرجاع أيقونة مناسبة لنوع الملف"""
        icons = {
            'pdf': '📕',
            'txt': '📝',
            'docx': '📘', 'doc': '📘',
            'pptx': '📙', 'ppt': '📙',
            'xlsx': '📗', 'xls': '📗', 'csv': '📗',
            'png': '🖼️', 'jpg': '🖼️', 'jpeg': '🖼️',
            'bmp': '🖼️', 'gif': '🖼️', 'webp': '🖼️',
            'mp3': '🎵', 'wav': '🎵', 'ogg': '🎵',
            'mp4': '🎬', 'avi': '🎬', 'mkv': '🎬',
        }
        return icons.get(extension, '📄')